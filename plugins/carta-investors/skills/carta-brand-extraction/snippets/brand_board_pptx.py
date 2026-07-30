"""Generate a brand board as a PPTX slide deck using python-pptx.

Dependencies: python-pptx, Pillow, requests.

Claude calls `generate_brand_board_pptx(data, output_path)` where `data` is the
classified brand data dict. Produces a multi-slide deck:
  Slide 1: Title slide (firm name + URL)
  Slide 2: Color Palette (swatches with hex codes)
  Slide 3: Typography (font specimens)
  Slide 4: Logo (on light and dark backgrounds)
  Slide 5: Imagery & Mood (image grid)
  Slide 6: Design Tokens summary
"""
from __future__ import annotations

import io
import tempfile

import requests
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _download_image(url: str) -> str | None:
    try:
        resp = requests.get(url, timeout=10, headers=_HEADERS)
        if not resp.ok:
            return None
        img = PILImage.open(io.BytesIO(resp.content))
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp, format="PNG")
        tmp.close()
        return tmp.name
    except Exception:
        return None


def _add_swatch(slide, left, top, width, height, hex_color: str) -> None:
    """Add a color swatch shape to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(hex_color)
    shape.line.fill.background()

    # Hex label below
    txBox = slide.shapes.add_textbox(left, top + height + Inches(0.05), width, Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = hex_color
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


def generate_brand_board_pptx(
    data: dict,
    output_path: str,
    firm_name: str = "",
) -> str:
    """Generate a brand board PPTX and return the output path.

    Args:
        data: Classified brand data dict (same shape as brand_board_pdf.py).
        output_path: Where to save the .pptx file.
        firm_name: Company name for the title slide.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    url = data.get("website_url", "")

    # ── Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = firm_name or "Brand Board"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Brand Board"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p2.alignment = PP_ALIGN.CENTER

    if url:
        p3 = tf.add_paragraph()
        p3.text = url
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p3.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Color Palette
    colors = data.get("colors", {})
    if any(colors.get(g) for g in ("primary", "secondary", "accent", "neutral")):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "COLOR PALETTE"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        row_y = Inches(1.2)
        for group_name in ("primary", "secondary", "accent", "neutral"):
            group = colors.get(group_name, [])
            if not group:
                continue

            # Group label
            txBox = slide.shapes.add_textbox(Inches(0.8), row_y - Inches(0.25), Inches(2), Inches(0.25))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = group_name.capitalize()
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

            for i, color_item in enumerate(group[:6]):
                _add_swatch(
                    slide,
                    left=Inches(0.8) + i * Inches(1.4),
                    top=row_y,
                    width=Inches(1.1),
                    height=Inches(1.1),
                    hex_color=color_item["hex"],
                )
            row_y += Inches(1.6)

    # ── Slide 3: Typography
    fonts = data.get("fonts", [])
    if fonts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "TYPOGRAPHY"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        y_pos = Inches(1.2)
        for font_item in fonts[:3]:
            family = font_item.get("family", "Unknown")
            weights = font_item.get("weights", ["400"])
            sizes = font_item.get("sizes", ["14px"])

            txBox = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11), Inches(1.8))
            tf = txBox.text_frame

            p = tf.paragraphs[0]
            p.text = family
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

            p2 = tf.add_paragraph()
            p2.text = "The quick brown fox jumps over the lazy dog"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

            p3 = tf.add_paragraph()
            p3.text = "ABCDEFGHIJKLMNOPQRSTUVWXYZ abcdefghijklmnopqrstuvwxyz 0123456789"
            p3.font.size = Pt(12)
            p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

            p4 = tf.add_paragraph()
            weights_str = ", ".join(str(w) for w in weights)
            sizes_str = ", ".join(str(s) for s in sizes)
            p4.text = f"Weights: {weights_str}  |  Sizes: {sizes_str}"
            p4.font.size = Pt(8)
            p4.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

            y_pos += Inches(2.0)

    # ── Slide 4: Logo
    logos = data.get("logos", [])
    if logos:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "LOGO"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        x_pos = Inches(1)
        for logo in logos[:3]:
            logo_url = logo.get("url")
            if not logo_url:
                continue
            tmp = _download_image(logo_url)
            if tmp:
                try:
                    slide.shapes.add_picture(
                        tmp, x_pos, Inches(2),
                        height=Inches(2),
                    )
                    x_pos += Inches(4)
                except Exception:
                    pass

    # ── Slide 5: Imagery & Mood
    imagery = data.get("imagery", [])
    if imagery:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "IMAGERY & MOOD"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        positions = [
            (Inches(0.8), Inches(1.2)), (Inches(4.8), Inches(1.2)),
            (Inches(8.8), Inches(1.2)), (Inches(0.8), Inches(4.2)),
            (Inches(4.8), Inches(4.2)), (Inches(8.8), Inches(4.2)),
        ]
        for i, img_item in enumerate(imagery[:6]):
            img_url = img_item.get("url") or img_item.get("src")
            if not img_url:
                continue
            tmp = _download_image(img_url)
            if tmp and i < len(positions):
                try:
                    x, y = positions[i]
                    slide.shapes.add_picture(tmp, x, y, width=Inches(3.6), height=Inches(2.5))
                except Exception:
                    pass

    # ── Slide 6: Design Tokens
    tokens = data.get("tokens", [])
    if tokens:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "DESIGN TOKENS"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11), Inches(5))
        tf = txBox.text_frame
        for token in tokens[:12]:
            p = tf.add_paragraph()
            p.text = f"{token.get('name', '')}: {token.get('value', '')}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    prs.save(output_path)
    return output_path
