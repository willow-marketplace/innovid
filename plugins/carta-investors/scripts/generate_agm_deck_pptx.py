"""
Generate a branded AGM (Annual General Meeting) presentation as PPTX.

Takes two JSON inputs:
  1. Brand data  — output of analyze_website.py (colors, typography, logos)
  2. Fund data   — fund metrics, portfolio, capital activity, etc.

Usage:
  uv run generate_agm_deck_pptx.py <brand.json> <fund_data.json> [output.pptx]

Or from Claude, call programmatically:
  generate_agm_deck_pptx(brand_data, fund_data, "output.pptx")

The deck adapts all colors, accents, and typography to whatever brand
identity was extracted — it is NOT hard-coded to any specific firm.
"""
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx"]
# ///
from __future__ import annotations

import json
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Brand theme resolution ───────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = h[0] * 2 + h[1] * 2 + h[2] * 2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Theme:
    """Resolved brand theme from analyze_website.py output."""

    def __init__(self, brand_data: dict) -> None:
        palette = brand_data.get("brand_palette", {})
        colors = brand_data.get("colors", [])

        self.firm_name = brand_data.get("firm_name", brand_data.get("title", ""))
        self.website_url = brand_data.get("url", brand_data.get("website_url", ""))

        self.primary = self._pick(palette, "primary", colors, "#1a1a1a")
        self.secondary = self._pick(palette, "secondary", colors, "#444444")
        self.accent = self._pick(palette, "accent", colors, "#0066cc")
        self.text_color = self._pick(palette, "text", colors, "#333333")
        self.bg_light = self._pick(palette, "background", colors, "#fafafa")
        self.neutral = self._pick(palette, "neutral", colors, "#222222")

        self.primary_rgb = _hex_to_rgb(self.primary)
        self.secondary_rgb = _hex_to_rgb(self.secondary)
        self.accent_rgb = _hex_to_rgb(self.accent)
        self.text_rgb = _hex_to_rgb(self.text_color)
        self.neutral_rgb = _hex_to_rgb(self.neutral)

        self.white = RGBColor(0xFF, 0xFF, 0xFF)
        self.gray_100 = RGBColor(0xFA, 0xFA, 0xFA)
        self.gray_300 = RGBColor(0xDD, 0xDD, 0xDD)
        self.gray_500 = RGBColor(0x97, 0x97, 0x97)
        self.gray_700 = RGBColor(0x33, 0x33, 0x33)
        self.gray_900 = RGBColor(0x22, 0x22, 0x22)

    @staticmethod
    def _pick(palette: dict, group: str, colors: list, fallback: str) -> str:
        items = palette.get(group, [])
        if items:
            return items[0].get("hex", fallback)
        for c in colors:
            if c.get("role") == group:
                return c.get("hex", fallback)
        return fallback


# ── Drawing helpers ──────────────────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text: str, *,
              font_size: int = 12, bold: bool = False, color: RGBColor | None = None,
              align=PP_ALIGN.LEFT, font_name: str = "Calibri") -> None:
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return txbox


def _add_shape_rect(slide, left, top, width, height, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_table(slide, t: Theme, left, top, width, headers: list[str],
               rows: list[list[str]]) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w = width // n_cols

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = tbl_shape.table

    for j, hdr in enumerate(headers):
        table.columns[j].width = col_w
        cell = table.cell(0, j)
        cell.text = hdr
        _style_cell(cell, font_size=9, bold=True, color=t.white, fill=t.secondary_rgb,
                     align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

    stripe = RGBColor(0xF7, 0xF9, 0xFC)
    for i, row in enumerate(rows):
        bg = stripe if (i % 2 == 1) else t.white
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            _style_cell(cell, font_size=9, bold=False, color=t.gray_900, fill=bg,
                         align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

    # Green line under header — set bottom border on header cells
    for j in range(n_cols):
        _set_cell_border(table.cell(0, j), bottom_color=t.primary_rgb, bottom_width=Pt(1.5))


def _style_cell(cell, font_size: int, bold: bool, color: RGBColor,
                fill: RGBColor, align=PP_ALIGN.LEFT) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.alignment = align
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


def _set_cell_border(cell, bottom_color: RGBColor | None = None,
                     bottom_width=Pt(1)) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    if bottom_color:
        ln = tcPr.find(qn("a:lnB"))
        if ln is None:
            ln = tcPr.makeelement(qn("a:lnB"), {})
            tcPr.append(ln)
        ln.set("w", str(int(bottom_width)))
        solid = ln.makeelement(qn("a:solidFill"), {})
        srgb = solid.makeelement(qn("a:srgbClr"), {"val": str(bottom_color)})
        solid.append(srgb)
        ln.append(solid)


def _add_kpi_card(slide, t: Theme, left, top, width, height,
                  value: str, label: str, accent: bool = False) -> None:
    bg = t.secondary_rgb if accent else t.gray_100
    _add_rounded_rect(slide, left, top, width, height, bg)

    val_color = t.white if accent else t.secondary_rgb
    lbl_color = RGBColor(0xCC, 0xCC, 0xCC) if accent else t.gray_500

    _add_text(slide, left + Inches(0.15), top + Inches(0.1),
              width - Inches(0.3), Inches(0.5), value,
              font_size=22, bold=True, color=val_color)
    _add_text(slide, left + Inches(0.15), top + height - Inches(0.35),
              width - Inches(0.3), Inches(0.3), label,
              font_size=9, color=lbl_color)


def _draw_top_bar(slide, t: Theme) -> None:
    _add_shape_rect(slide, 0, 0, SLIDE_W, Inches(0.08), t.secondary_rgb)
    _add_shape_rect(slide, 0, 0, Emu(int(SLIDE_W * 0.3)), Inches(0.08), t.primary_rgb)


def _draw_footer(slide, t: Theme, page_num: int, total: int, firm_name: str,
                 dark: bool = False) -> None:
    color = RGBColor(0xAA, 0xAA, 0xAA) if dark else t.gray_500
    _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
              f"{firm_name}  |  Annual General Meeting  |  Confidential",
              font_size=7, color=color)
    _add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.3),
              f"{page_num} / {total}", font_size=7, color=color, align=PP_ALIGN.RIGHT)


# ── Slide builders ───────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.secondary_rgb)

    # Accent bar
    _add_shape_rect(slide, Inches(0.6), Inches(2.8), Inches(0.8), Inches(0.06), t.primary_rgb)

    _add_text(slide, Inches(0.6), Inches(3.0), Inches(10), Inches(1),
              "Annual General Meeting", font_size=42, bold=True, color=t.white)

    fund_name = fd.get("fund_name", "Fund")
    _add_text(slide, Inches(0.6), Inches(3.85), Inches(10), Inches(0.5),
              fund_name, font_size=18, color=t.primary_rgb)

    as_of = fd.get("as_of_date", "")
    subtitle = f"{as_of}  |  Confidential" if as_of else "Confidential"
    _add_text(slide, Inches(0.6), Inches(4.4), Inches(10), Inches(0.4),
              subtitle, font_size=13, color=RGBColor(0x99, 0x99, 0x99))

    firm = t.firm_name or fd.get("firm_name", "")
    _add_text(slide, Inches(0.6), SLIDE_H - Inches(1), Inches(4), Inches(0.4),
              firm.upper(), font_size=14, bold=True, color=RGBColor(0x55, 0x55, 0x55))

    _draw_footer(slide, t, pg, total, firm, dark=True)


def _slide_agenda(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    _add_text(slide, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
              "OVERVIEW", font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(6), Inches(0.6),
              "Agenda", font_size=28, bold=True, color=t.secondary_rgb)

    agenda = fd.get("agenda", list(_ACTIVE_LABELS))

    y = Inches(1.65)
    for i, item in enumerate(agenda, 1):
        _add_text(slide, Inches(0.6), y, Inches(0.6), Inches(0.4),
                  f"{i:02d}", font_size=22, bold=True, color=t.primary_rgb)
        _add_text(slide, Inches(1.3), y + Inches(0.03), Inches(8), Inches(0.4),
                  item, font_size=14, color=t.gray_900)
        # Divider line
        _add_shape_rect(slide, Inches(0.6), y + Inches(0.45), SLIDE_W - Inches(1.2),
                        Inches(0.01), t.gray_300)
        y += Inches(0.6)

    _draw_footer(slide, t, pg, total, firm)


def _slide_fund_overview(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_fund_overview, 0):02d}  {SLIDE_LABELS[_slide_fund_overview].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Fund Overview", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("fund_overview_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    summary = fd.get("fund_overview_summary", "")
    if summary:
        _add_text(slide, Inches(0.6), y_after_kpis, Inches(11), Inches(0.6),
                  summary, font_size=11, color=t.gray_700)

    _draw_footer(slide, t, pg, total, firm)


def _slide_fund_performance(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_fund_performance, 0):02d}  {SLIDE_LABELS[_slide_fund_performance].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Fund Performance Summary", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    perf = fd.get("performance_table", {})
    if perf:
        headers = perf.get("headers", [])
        rows = perf.get("rows", [])
        table_w = SLIDE_W - Inches(1.2)
        _add_table(slide, t, Inches(0.6), y_after_kpis, table_w, headers, rows)

    disclaimer = fd.get("performance_disclaimer", "")
    if disclaimer:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  disclaimer, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_benchmarks(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_benchmarks, 0):02d}  {SLIDE_LABELS[_slide_benchmarks].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Performance Benchmarks", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("benchmark_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    tbl = fd.get("benchmark_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("benchmark_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_nav_bridge(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_nav_bridge, 0):02d}  {SLIDE_LABELS[_slide_nav_bridge].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "NAV Bridge", font_size=24, bold=True, color=t.secondary_rgb)

    tbl = fd.get("nav_bridge_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("nav_bridge_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_portfolio(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_portfolio, 0):02d}  {SLIDE_LABELS[_slide_portfolio].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Portfolio Companies", font_size=24, bold=True, color=t.secondary_rgb)

    portfolio = fd.get("portfolio_table", {})
    if portfolio:
        headers = portfolio.get("headers", [])
        rows = portfolio.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("portfolio_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_sector_allocation(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_sector_allocation, 0):02d}  {SLIDE_LABELS[_slide_sector_allocation].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Sector Allocation", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("sector_allocation_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    tbl = fd.get("sector_allocation_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    _draw_footer(slide, t, pg, total, firm)


def _slide_vintage_allocation(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_vintage_allocation, 0):02d}  {SLIDE_LABELS[_slide_vintage_allocation].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Vintage Allocation", font_size=24, bold=True, color=t.secondary_rgb)

    tbl = fd.get("vintage_allocation_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("vintage_allocation_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_top_investments(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_top_investments, 0):02d}  {SLIDE_LABELS[_slide_top_investments].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Top Performing Investments", font_size=24, bold=True, color=t.secondary_rgb)

    top_table = fd.get("top_investments_table", {})
    if top_table:
        headers = top_table.get("headers", [])
        rows = top_table.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("top_investments_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_realized(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_realized, 0):02d}  {SLIDE_LABELS[_slide_realized].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Realized Investments", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("realized_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    tbl = fd.get("realized_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("realized_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_watchlist(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_watchlist, 0):02d}  {SLIDE_LABELS[_slide_watchlist].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Watchlist", font_size=24, bold=True, color=t.secondary_rgb)

    tbl = fd.get("watchlist_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("watchlist_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_capital_activity(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_capital_activity, 0):02d}  {SLIDE_LABELS[_slide_capital_activity].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Capital Activity & Cash Flows", font_size=24, bold=True, color=t.secondary_rgb)

    cap_kpis = fd.get("capital_kpis", [])
    y_after_kpis = Inches(1.55)
    if cap_kpis:
        n = len(cap_kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(cap_kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 1))
        y_after_kpis = Inches(2.65)

    cap_table = fd.get("capital_table", {})
    if cap_table:
        headers = cap_table.get("headers", [])
        rows = cap_table.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("capital_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_fees(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_fees, 0):02d}  {SLIDE_LABELS[_slide_fees].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Fees & Expenses", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("fees_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    tbl = fd.get("fees_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("fees_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_balance_sheet(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_balance_sheet, 0):02d}  {SLIDE_LABELS[_slide_balance_sheet].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Balance Sheet", font_size=24, bold=True, color=t.secondary_rgb)

    tbl = fd.get("balance_sheet_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), Inches(1.55), SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("balance_sheet_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_lp_summary(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_lp_summary, 0):02d}  {SLIDE_LABELS[_slide_lp_summary].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "LP Base Summary", font_size=24, bold=True, color=t.secondary_rgb)

    kpis = fd.get("lp_summary_kpis", [])
    y_after_kpis = Inches(1.55)
    if kpis:
        n = len(kpis)
        card_w = (SLIDE_W - Inches(1.2) - Inches(0.2) * (n - 1)) // n
        card_h = Inches(0.9)
        for i, kpi in enumerate(kpis):
            left = Inches(0.6) + i * (card_w + Inches(0.2))
            _add_kpi_card(slide, t, left, Inches(1.55), card_w, card_h,
                          kpi["value"], kpi["label"],
                          accent=kpi.get("accent", i % 2 == 0))
        y_after_kpis = Inches(2.65)

    tbl = fd.get("lp_summary_table", {})
    if tbl:
        headers = tbl.get("headers", [])
        rows = tbl.get("rows", [])
        _add_table(slide, t, Inches(0.6), y_after_kpis, SLIDE_W - Inches(1.2), headers, rows)

    note = fd.get("lp_summary_note", "")
    if note:
        _add_text(slide, Inches(0.6), SLIDE_H - Inches(0.75), Inches(11), Inches(0.3),
                  note, font_size=7, color=t.gray_500)

    _draw_footer(slide, t, pg, total, firm)


def _slide_esg(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_esg, 0):02d}  {SLIDE_LABELS[_slide_esg].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "ESG & Impact", font_size=24, bold=True, color=t.secondary_rgb)

    themes = fd.get("esg_themes", [])
    y = Inches(1.6)
    for item in themes:
        title = item.get("title", "")
        desc = item.get("description", "")

        circle = slide.shapes.add_shape(
            9, Inches(0.6), y + Inches(0.02), Inches(0.13), Inches(0.13))
        circle.fill.solid()
        circle.fill.fore_color.rgb = t.primary_rgb
        circle.line.fill.background()

        _add_text(slide, Inches(0.9), y - Inches(0.03), Inches(10), Inches(0.35),
                  title, font_size=13, bold=True, color=t.secondary_rgb)

        _add_text(slide, Inches(0.9), y + Inches(0.3), Inches(11), Inches(0.6),
                  desc, font_size=10, color=t.gray_700)

        y += Inches(1.2)

    _draw_footer(slide, t, pg, total, firm)


def _slide_outlook(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_outlook, 0):02d}  {SLIDE_LABELS[_slide_outlook].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Market Outlook & Strategy", font_size=24, bold=True, color=t.secondary_rgb)

    themes = fd.get("outlook_themes", [])
    y = Inches(1.6)
    for item in themes:
        title = item.get("title", "")
        desc = item.get("description", "")

        circle = slide.shapes.add_shape(
            9, Inches(0.6), y + Inches(0.02), Inches(0.13), Inches(0.13))
        circle.fill.solid()
        circle.fill.fore_color.rgb = t.primary_rgb
        circle.line.fill.background()

        _add_text(slide, Inches(0.9), y - Inches(0.03), Inches(10), Inches(0.35),
                  title, font_size=13, bold=True, color=t.secondary_rgb)

        _add_text(slide, Inches(0.9), y + Inches(0.3), Inches(11), Inches(0.6),
                  desc, font_size=10, color=t.gray_700)

        y += Inches(1.2)

    _draw_footer(slide, t, pg, total, firm)


def _slide_risks(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.white)
    _draw_top_bar(slide, t)
    firm = t.firm_name or fd.get("firm_name", "")

    sec = f"{_SECTION_NUMS.get(_slide_risks, 0):02d}  {SLIDE_LABELS[_slide_risks].upper()}"
    _add_text(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.3),
              sec, font_size=8, bold=True, color=t.primary_rgb)
    _add_text(slide, Inches(0.6), Inches(0.75), Inches(8), Inches(0.6),
              "Key Risks", font_size=24, bold=True, color=t.secondary_rgb)

    themes = fd.get("risk_themes", [])
    y = Inches(1.6)
    for item in themes:
        title = item.get("title", "")
        desc = item.get("description", "")

        circle = slide.shapes.add_shape(
            9, Inches(0.6), y + Inches(0.02), Inches(0.13), Inches(0.13))
        circle.fill.solid()
        circle.fill.fore_color.rgb = t.primary_rgb
        circle.line.fill.background()

        _add_text(slide, Inches(0.9), y - Inches(0.03), Inches(10), Inches(0.35),
                  title, font_size=13, bold=True, color=t.secondary_rgb)

        _add_text(slide, Inches(0.9), y + Inches(0.3), Inches(11), Inches(0.6),
                  desc, font_size=10, color=t.gray_700)

        y += Inches(1.2)

    _draw_footer(slide, t, pg, total, firm)


def _slide_closing(prs: Presentation, t: Theme, fd: dict, pg: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, t.secondary_rgb)
    firm = t.firm_name or fd.get("firm_name", "")

    _add_shape_rect(slide, Inches(0.6), Inches(2.9), Inches(0.8), Inches(0.06), t.primary_rgb)

    _add_text(slide, Inches(0.6), Inches(3.1), Inches(8), Inches(0.8),
              "Thank You", font_size=36, bold=True, color=t.white)
    _add_text(slide, Inches(0.6), Inches(3.8), Inches(8), Inches(0.5),
              "Questions & Discussion", font_size=16,
              color=RGBColor(0xAA, 0xAA, 0xAA))

    contact = fd.get("contact", {})
    contact_lines = contact.get("lines", [firm])
    y = Inches(4.6)
    for line in contact_lines:
        _add_text(slide, Inches(0.6), y, Inches(6), Inches(0.3),
                  line, font_size=12, color=RGBColor(0x88, 0x88, 0x88))
        y += Inches(0.3)

    _add_text(slide, Inches(0.6), SLIDE_H - Inches(1), Inches(4), Inches(0.4),
              firm.upper(), font_size=14, bold=True, color=RGBColor(0x55, 0x55, 0x55))

    _draw_footer(slide, t, pg, total, firm, dark=True)


# ── Main entry point ─────────────────────────────────────────────────────────

SLIDE_BUILDERS = [
    _slide_cover,
    _slide_agenda,
    _slide_fund_overview,
    _slide_fund_performance,
    _slide_benchmarks,
    _slide_nav_bridge,
    _slide_portfolio,
    _slide_sector_allocation,
    _slide_vintage_allocation,
    _slide_top_investments,
    _slide_realized,
    _slide_watchlist,
    _slide_capital_activity,
    _slide_fees,
    _slide_balance_sheet,
    _slide_lp_summary,
    _slide_esg,
    _slide_outlook,
    _slide_risks,
    _slide_closing,
]

SLIDE_LABELS = {
    _slide_fund_overview: "Fund Overview",
    _slide_fund_performance: "Fund Performance",
    _slide_benchmarks: "Performance Benchmarks",
    _slide_nav_bridge: "NAV Bridge",
    _slide_portfolio: "Portfolio Overview",
    _slide_sector_allocation: "Sector Allocation",
    _slide_vintage_allocation: "Vintage Allocation",
    _slide_top_investments: "Top Investments",
    _slide_realized: "Realized Investments",
    _slide_watchlist: "Watchlist",
    _slide_capital_activity: "Capital Activity",
    _slide_fees: "Fees & Expenses",
    _slide_balance_sheet: "Balance Sheet",
    _slide_lp_summary: "LP Base Summary",
    _slide_esg: "ESG & Impact",
    _slide_outlook: "Market Outlook",
    _slide_risks: "Key Risks",
}

_NON_CONTENT = {_slide_cover, _slide_agenda, _slide_closing}

_SECTION_NUMS: dict = {}
_ACTIVE_LABELS: list[str] = []


def generate_agm_deck_pptx(brand_data: dict, fund_data: dict, output_path: str) -> str:
    """Generate a branded AGM deck as PPTX.

    :param brand_data: Output of analyze_website.py (or manually constructed brand dict).
    :param fund_data: Fund metrics, portfolio, capital activity, outlook themes, etc.
    :param output_path: Where to write the PPTX.
    :returns: The output path.
    """
    global _SECTION_NUMS, _ACTIVE_LABELS

    t = Theme(brand_data)

    data_requirements = {
        _slide_fund_overview: ("fund_overview_kpis",),
        _slide_fund_performance: ("kpis", "performance_table"),
        _slide_benchmarks: ("benchmark_kpis", "benchmark_table"),
        _slide_nav_bridge: ("nav_bridge_table",),
        _slide_portfolio: ("portfolio_table",),
        _slide_sector_allocation: ("sector_allocation_kpis", "sector_allocation_table"),
        _slide_vintage_allocation: ("vintage_allocation_table",),
        _slide_top_investments: ("top_investments_table",),
        _slide_realized: ("realized_kpis", "realized_table"),
        _slide_watchlist: ("watchlist_table",),
        _slide_capital_activity: ("capital_kpis", "capital_table"),
        _slide_fees: ("fees_kpis", "fees_table"),
        _slide_balance_sheet: ("balance_sheet_table",),
        _slide_lp_summary: ("lp_summary_kpis", "lp_summary_table"),
        _slide_esg: ("esg_themes",),
        _slide_outlook: ("outlook_themes",),
        _slide_risks: ("risk_themes",),
    }

    active: list = []
    for slide_fn in SLIDE_BUILDERS:
        reqs = data_requirements.get(slide_fn)
        if reqs and not any(fund_data.get(k) for k in reqs):
            continue
        active.append(slide_fn)

    _SECTION_NUMS = {}
    _ACTIVE_LABELS = []
    sec_num = 0
    for slide_fn in active:
        if slide_fn in _NON_CONTENT:
            continue
        sec_num += 1
        _SECTION_NUMS[slide_fn] = sec_num
        _ACTIVE_LABELS.append(SLIDE_LABELS[slide_fn])

    total = len(active)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_fn in enumerate(active, 1):
        slide_fn(prs, t, fund_data, i, total)

    prs.save(output_path)
    return output_path


def main() -> None:
    if len(sys.argv) < 3:
        print("Usage: generate_agm_deck_pptx.py <brand.json> <fund_data.json> [output.pptx]")
        sys.exit(1)

    with open(sys.argv[1]) as f:
        brand_data = json.load(f)
    with open(sys.argv[2]) as f:
        fund_data = json.load(f)

    output = sys.argv[3] if len(sys.argv) > 3 else "agm-deck.pptx"
    result = generate_agm_deck_pptx(brand_data, fund_data, output)
    print(f"Generated: {result}")


if __name__ == "__main__":
    main()
